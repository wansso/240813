import streamlit as st
from pptx import Presentation
from transformers import pipeline
import os

# Streamlit 앱 제목
st.title("PPT 요약기💬")

# 처음 1번만 실행하기 위한 코드
if "messages" not in st.session_state:
    st.session_state["messages"] = []

if "ppt_chain" not in st.session_state:
    st.session_state["ppt_chain"] = None

if "ppt_retriever" not in st.session_state:
    st.session_state["ppt_retriever"] = None

# 사이드바 생성
with st.sidebar:
    # 초기화 버튼 생성
    clear_btn = st.button("대화 초기화")

    # 파일 업로드
    uploaded_file = st.file_uploader("파일 업로드", type=["pptx"])

    # 모델 선택 메뉴
    selected_model = st.selectbox(
        "요약 모델 선택", ["distilbart-cnn-12-6", "facebook/bart-large-cnn"], index=0
    )

    update_btn = st.button("설정 업데이트")


# 이전 대화를 출력
def print_messages():
    for chat_message in st.session_state["messages"]:
        st.chat_message(chat_message['role']).write(chat_message['content'])


# 새로운 메시지를 추가
def add_message(role, message):
    st.session_state["messages"].append({"role": role, "content": message})


# PPT 파일에서 텍스트를 추출하는 함수
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_runs.append("\n".join(slide_text))

    return text_runs


# 텍스트 요약 함수
@st.cache_resource(show_spinner="업로드한 파일을 처리 중입니다...")
def summarize_text(text_data, model_name="distilbart-cnn-12-6"):
    summarizer = pipeline("summarization", model=model_name)
    summaries = []
    for text in text_data:
        summary = summarizer(text, max_length=150, min_length=40, do_sample=False)[0]['summary_text']
        summaries.append(summary)
    return summaries


# 파일이 업로드 되었을 때
if uploaded_file:
    # 파일에서 텍스트를 추출하고 요약 생성 (작업시간이 오래 걸릴 예정...)
    ppt_text = extract_text_from_pptx(uploaded_file)
    summarized_text = summarize_text(ppt_text, model_name=selected_model)
    st.session_state["ppt_chain"] = summarized_text

# 초기화 버튼이 눌리면...
if clear_btn:
    st.session_state["messages"] = []

if update_btn:
    if st.session_state["ppt_chain"] is not None:
        summarized_text = summarize_text(
            extract_text_from_pptx(uploaded_file),
            model_name=selected_model
        )
        st.session_state["ppt_chain"] = summarized_text
    else:
        ppt_text = extract_text_from_pptx(uploaded_file)
        summarized_text = summarize_text(ppt_text, model_name=selected_model)
        st.session_state["ppt_chain"] = summarized_text

# 이전 대화 기록 출력
print_messages()

# 요약된 내용을 출력
if st.session_state["ppt_chain"]:
    for idx, summary in enumerate(st.session_state["ppt_chain"]):
        st.write(f"Slide {idx + 1} Summary: {summary}")

# 사용자 입력을 받아서 대화에 추가
user_input = st.chat_input("궁금한 내용을 물어보세요!")

# 경고 메시지를 띄우기 위한 빈 영역
warning_msg = st.empty()

# 만약에 사용자 입력이 들어오면...
if user_input:
    # 요약된 내용이 있는지 확인
    if st.session_state["ppt_chain"] is not None:
        st.chat_message("user").write(user_input)
        
        # 간단한 응답 처리 (예: 특정 슬라이드에 대한 질문)
        response = f"Slide {user_input}에 대한 요약: {st.session_state['ppt_chain'][int(user_input) - 1]}"
        
        with st.chat_message("assistant"):
            st.write(response)

        # 대화기록을 저장한다.
        add_message("user", user_input)
        add_message("assistant", response)
    else:
        # 파일을 업로드 하라는 경고 메시지 출력
        warning_msg.error("파일을 업로드 해주세요.")

